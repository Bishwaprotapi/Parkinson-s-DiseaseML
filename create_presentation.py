from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_title_slide(prs):
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Parkinson's Disease Detection\nusing Machine Learning"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    
    subtitle.text = "Machine Learning Project\n" + "Date: " + "2024"

def create_intro_slide(prs):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Introduction"
    
    tf = body_shape.text_frame
    tf.text = "Parkinson's Disease Overview"
    
    p = tf.add_paragraph()
    p.text = "• Progressive neurological disorder"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Affects movement and coordination"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Early detection is crucial for better management"
    p.level = 1

def create_project_overview(prs):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Project Overview"
    
    tf = body_shape.text_frame
    tf.text = "Project Goals"
    
    p = tf.add_paragraph()
    p.text = "• Develop ML model for Parkinson's detection"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Analyze key features and patterns"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Achieve high accuracy in prediction"
    p.level = 1

def create_methodology(prs):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Methodology"
    
    tf = body_shape.text_frame
    tf.text = "Project Approach"
    
    p = tf.add_paragraph()
    p.text = "• Data Collection and Preprocessing"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Feature Selection and Engineering"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Model Training and Evaluation"
    p.level = 1

def create_results(prs):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Results"
    
    tf = body_shape.text_frame
    tf.text = "Key Findings"
    
    p = tf.add_paragraph()
    p.text = "• Model Performance Metrics"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Accuracy and Precision Scores"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Feature Importance Analysis"
    p.level = 1

def create_conclusion(prs):
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = "Conclusion"
    
    tf = body_shape.text_frame
    tf.text = "Summary"
    
    p = tf.add_paragraph()
    p.text = "• Successful implementation of ML model"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Potential for real-world application"
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = "• Future improvements and research directions"
    p.level = 1

def main():
    prs = Presentation()
    
    # Create slides
    create_title_slide(prs)
    create_intro_slide(prs)
    create_project_overview(prs)
    create_methodology(prs)
    create_results(prs)
    create_conclusion(prs)
    
    # Save the presentation
    prs.save('Parkinson_Disease_ML_Presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    main() 